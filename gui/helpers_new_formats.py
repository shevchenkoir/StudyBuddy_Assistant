# helpers_new_formats.py
from PIL import Image
from io import BytesIO
import pandas as pd

def extract_text_from_pptx(file_path: str) -> str:
    """
    Извлечение текста из PPTX.
    Требует: python-pptx
    """
    from pptx import Presentation
    prs = Presentation(file_path)
    lines = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                lines.append(shape.text)
    return "\n".join(lines)

def extract_text_from_xlsx(file_path: str, max_cells: int = 20000) -> str:
    """
    Преобразует листы Excel в плоский текст.
    Требует: pandas, openpyxl
    Ограничивает количество ячеек, чтобы не взорваться по памяти.
    """
    xl = pd.ExcelFile(file_path)
    parts = []
    cells_seen = 0
    for sheet in xl.sheet_names:
        df = xl.parse(sheet)
        for row in df.astype(str).values.tolist():
            parts.append(" ".join(row))
            cells_seen += len(row)
            if cells_seen >= max_cells:
                parts.append("\n[...обрезано для индексации...]")
                return "\n".join(parts)
    return "\n".join(parts)

def extract_text_from_csv(file_path: str, encoding: str = "utf-8") -> str:
    """
    Читает CSV как обычный текст.
    """
    with open(file_path, "r", encoding=encoding, errors="ignore") as f:
        return f.read()

def extract_text_from_image(file_path: str, lang: str = "eng") -> str:
    """
    OCR изображений через Tesseract.
    Требует: pillow, pytesseract (+ установленный системно tesseract-ocr)
    lang: 'eng', 'rus' или 'eng+rus'
    """
    import pytesseract
    image = Image.open(file_path)
    return pytesseract.image_to_string(image, lang=lang) or ""